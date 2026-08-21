"""Tests for document text extraction and validation helpers (PDF, DOCX, PPTX)."""

from __future__ import annotations

from io import BytesIO

import docx
import pptx
from pypdf import PdfWriter

from pdf_text import (
    MAX_UPLOAD_SIZE_BYTES,
    extract_document_text,
    extract_docx_text,
    extract_pptx_text,
    pdf_is_likely_image_only,
)


def test_pdf_is_likely_image_only_rejects_blank_pdf():
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = BytesIO()
    writer.write(buf)
    assert pdf_is_likely_image_only(buf.getvalue()) is True


def test_pdf_is_likely_image_only_accepts_text_pdf(monkeypatch):
    long_text = "Bilkeys extractable chapter text. " * 30

    class FakePage:
        def extract_text(self):
            return long_text

    class FakeReader:
        pages = [FakePage()]

    monkeypatch.setattr("pdf_text.PdfReader", lambda _: FakeReader())
    assert pdf_is_likely_image_only(b"fake-pdf") is False


def test_docx_text_extraction():
    doc = docx.Document()
    doc.add_heading("Chapter 1: Introduction to Biology", level=1)
    doc.add_paragraph("Cells are the basic structural and functional units of life.")
    buf = BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    text = extract_docx_text(docx_bytes)
    assert "Chapter 1: Introduction to Biology" in text
    assert "Cells are the basic structural" in text
    assert pdf_is_likely_image_only(docx_bytes, filename="test.docx") is False


def test_pptx_text_extraction():
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Presentation Title: Organic Chemistry"
    subtitle = slide.placeholders[1]
    subtitle.text = "Chapter 2: Alkanes and Cycloalkanes"
    buf = BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    text = extract_pptx_text(pptx_bytes)
    assert "Presentation Title: Organic Chemistry" in text
    assert "Chapter 2: Alkanes and Cycloalkanes" in text
    assert pdf_is_likely_image_only(pptx_bytes, filename="lecture.pptx") is False


def test_extract_document_text_routes_by_extension():
    doc = docx.Document()
    doc.add_paragraph("Sample Word Content")
    buf = BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    text = extract_document_text(docx_bytes, filename="document.docx")
    assert "Sample Word Content" in text


def test_max_upload_size_constant():
    assert MAX_UPLOAD_SIZE_BYTES == 20 * 1024 * 1024
