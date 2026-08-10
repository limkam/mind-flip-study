"""Shared PDF text extraction."""

from __future__ import annotations

import re
from io import BytesIO

from pypdf import PdfReader

TOC_SAMPLE_CHARS = 100_000
TOC_SAMPLE_MAX = 600_000

_HEADING_LINE = re.compile(
    r"^(?:\d+\.?\s+|(?:chapter|ch\.?|section|appendix|part)\s+)",
    re.I,
)


def extract_pdf_text(
    data: bytes,
    *,
    max_pages: int | None = None,
    include_page_markers: bool = False,
) -> str:
    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    pages = reader.pages if max_pages is None else reader.pages[:max_pages]
    for page_number, page in enumerate(pages, start=1):
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t:
            if include_page_markers:
                parts.append(f"[PAGE {page_number}]")
            parts.append(t)
    return "\n".join(parts)


# Minimum average extractable characters per page for a text-based PDF.
_MIN_CHARS_PER_PAGE = 40
# Absolute floor — very short docs may be valid.
_MIN_TOTAL_TEXT_CHARS = 80

# Pages to scan for image-only detection (full book is not needed).
_MAX_IMAGE_CHECK_PAGES = 12

IMAGE_ONLY_PDF_MESSAGE = (
    "This document appears to be scanned photos or an image-only document. "
    "MindFlip cannot process image documents yet — please upload a text-based document "
    "(for example, standard PDF, Word .docx, PowerPoint .pptx, or Google Docs export)."
)


MAX_UPLOAD_SIZE_BYTES = 10 * 1024 * 1024  # 10 MB

def extract_docx_text(data: bytes, *, include_page_markers: bool = False) -> str:
    try:
        import docx
    except ImportError:
        return ""
    doc = docx.Document(BytesIO(data))
    paragraphs: list[str] = []
    section_count = 1
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            if include_page_markers and p.style and p.style.name and "Heading 1" in p.style.name:
                paragraphs.append(f"[PAGE {section_count}]")
                section_count += 1
            paragraphs.append(t)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def extract_pptx_text(data: bytes, *, include_page_markers: bool = False) -> str:
    try:
        import pptx
    except ImportError:
        return ""
    prs = pptx.Presentation(BytesIO(data))
    parts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            if include_page_markers:
                parts.append(f"[PAGE {slide_idx}]")
            parts.extend(slide_texts)
    return "\n".join(parts)


def extract_document_text(
    data: bytes,
    filename: str = "",
    *,
    include_page_markers: bool = False,
) -> str:
    fn = filename.lower()
    if fn.endswith(".docx"):
        return extract_docx_text(data, include_page_markers=include_page_markers)
    if fn.endswith(".pptx"):
        return extract_pptx_text(data, include_page_markers=include_page_markers)
    return extract_pdf_text(data, include_page_markers=include_page_markers)


def pdf_is_likely_image_only(data: bytes, filename: str = "", *, max_pages: int = _MAX_IMAGE_CHECK_PAGES) -> bool:
    """
    Heuristic: PDFs that are photos/scans have little or no extractable text.
    DOCX and PPTX with extractable text pass automatically.
    """
    fn = filename.lower()
    if fn.endswith(".docx"):
        text = extract_docx_text(data)
        return len(text.strip()) < 10
    if fn.endswith(".pptx"):
        text = extract_pptx_text(data)
        return len(text.strip()) < 10

    # Auto-detect ZIP archives (DOCX / PPTX) if filename extension was not provided
    if data[:4] == b"PK\x03\x04":
        doc_text = extract_docx_text(data)
        if doc_text.strip():
            return len(doc_text.strip()) < 10
        ppt_text = extract_pptx_text(data)
        if ppt_text.strip():
            return len(ppt_text.strip()) < 10

    try:
        reader = PdfReader(BytesIO(data))
    except Exception:
        doc_text = extract_docx_text(data)
        if doc_text.strip():
            return len(doc_text.strip()) < 10
        ppt_text = extract_pptx_text(data)
        if ppt_text.strip():
            return len(ppt_text.strip()) < 10
        return True

    pages = reader.pages
    if not pages:
        return True

    sample = pages[: max(1, min(max_pages, len(pages)))]
    total_chars = 0
    sparse_pages = 0
    for i, page in enumerate(sample):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        char_count = len(text)
        total_chars += char_count
        if char_count < _MIN_CHARS_PER_PAGE:
            sparse_pages += 1
        # Clearly text-based after a couple of pages — skip scanning the rest.
        if i >= 1 and total_chars >= 400:
            return False

    if total_chars < _MIN_TOTAL_TEXT_CHARS:
        return True

    avg_chars = total_chars / len(sample)
    if avg_chars < _MIN_CHARS_PER_PAGE:
        return True

    if sparse_pages >= max(1, int(len(sample) * 0.85)):
        return True

    return False


def _collect_heading_lines(full_text: str, *, limit: int = 500) -> list[str]:
    """Scan the full document for section-heading lines to give AI document-wide structure."""
    found: list[str] = []
    seen: set[str] = set()
    for line in full_text.splitlines():
        stripped = re.sub(r"\s+", " ", line.strip())
        if len(stripped) < 4 or len(stripped) > 150:
            continue
        if not _HEADING_LINE.match(stripped):
            continue
        key = stripped.lower()
        if key in seen:
            continue
        seen.add(key)
        found.append(stripped[:150])
        if len(found) >= limit:
            break
    return found


def toc_sample_text(full_text: str) -> str:
    """
    Return the marked document text within the model context budget.
    """
    text = full_text.strip()
    if not text:
        return ""

    return text[:TOC_SAMPLE_MAX]
